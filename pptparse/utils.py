import asyncio
import json
import logging
import os
import shutil
import subprocess
import tempfile
import traceback
from shutil import which
from time import sleep, time
from typing import Any, Optional

from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.parts.image import Image
from pptx.shapes.group import GroupShape
from pptx.text.text import _Paragraph, _Run
from pptx.util import Length, Pt
from tenacity import RetryCallState, retry, stop_after_attempt, wait_fixed

def get_logger(name="pptagent", level=None):
    """
    Get a logger with the specified name and level.

    Args:
        name (str): The name of the logger.
        level (int): The logging level (default: logging.INFO).

    Returns:
        logging.Logger: A configured logger instance.
    """
    if level is None:
        level = int(os.environ.get("LOG_LEVEL", logging.INFO))

    logger = logging.getLogger(name)
    logger.setLevel(level)

    # Check if the logger already has handlers to avoid duplicates
    if not logger.handlers:
        # Create console handler and set level
        console_handler = logging.StreamHandler()
        console_handler.setLevel(level)

        # Create formatter
        formatter = logging.Formatter(
            "%(asctime)s - %(levelname)s - %(filename)s:%(lineno)d - %(message)s"
        )
        console_handler.setFormatter(formatter)

        # Add handler to logger
        logger.addHandler(console_handler)

    return logger


logger = get_logger(__name__)

if which("soffice") is None:
    logging.warning("soffice is not installed, pptx to images conversion will not work")

# Set of supported image extensions
IMAGE_EXTENSIONS: set[str] = {
    "bmp",
    "jpg",
    "jpeg",
    "pgm",
    "png",
    "ppm",
    "tif",
    "tiff",
    "webp",
}

# Common colors and measurements
BLACK = RGBColor(0, 0, 0)
YELLOW = RGBColor(255, 255, 0)
BLUE = RGBColor(0, 0, 255)
BORDER_LEN = Pt(2)
BORDER_OFFSET = Pt(2)
LABEL_LEN = Pt(24)
FONT_LEN = Pt(20)


def is_image_path(file: str) -> bool:
    """
    Check if a file path is an image based on its extension.

    Args:
        file (str): The file path to check.

    Returns:
        bool: True if the file is an image, False otherwise.
    """
    return file.split(".")[-1].lower() in IMAGE_EXTENSIONS


def runs_merge(paragraph: _Paragraph) -> Optional[_Run]:
    """
    Merge all runs in a paragraph into a single run.

    Args:
        paragraph (_Paragraph): The paragraph to merge runs in.

    Returns:
        Optional[_Run]: The merged run, or None if there are no runs.
    """
    runs = paragraph.runs

    # Handle field codes
    if len(runs) == 0:
        runs = [
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        ]
    if len(runs) == 1:
        return runs[0]
    if len(runs) == 0:
        return None

    # Find the run with the most text
    run = max(runs, key=lambda x: len(x.text))
    run.text = paragraph.text

    # Remove other runs
    for r in runs:
        if r != run:
            r._r.getparent().remove(r._r)
    return run


def older_than(filepath: str, seconds: int = 10, wait: bool = False) -> bool:
    """
    Check if a file is older than a specified number of seconds.

    Args:
        filepath (str): The path to the file.
        seconds (int): The number of seconds to check against.
        wait (bool): Whether to wait for the file to exist.

    Returns:
        bool: True if the file is older than the specified number of seconds, False otherwise.
    """
    if not os.path.exists(filepath):
        while wait:
            logger.info("waiting for: %s", filepath)
            sleep(1)
            if os.path.exists(filepath):
                sleep(seconds)
                return True
        return False
    file_creation_time = os.path.getctime(filepath)
    current_time = time()
    return seconds < (current_time - file_creation_time)


def tenacity_log(retry_state: RetryCallState) -> None:
    """
    Log function for tenacity retries.

    Args:
        retry_state (RetryCallState): The retry state.
    """
    logger.warning("tenacity retry: %s", retry_state)
    traceback.print_tb(retry_state.outcome.exception().__traceback__)


# Create a tenacity decorator with custom settings
def tenacity_decorator(_func=None, *, wait: int = 3, stop: int = 5):
    def decorator(func):
        return retry(wait=wait_fixed(wait), stop=stop_after_attempt(stop))(func)

    if _func is None:
        # Called with arguments
        return decorator
    else:
        # Called without arguments
        return decorator(_func)


TABLE_CSS = """
table {
    border-collapse: collapse;  /* Merge borders */
    width: auto;               /* Width adapts to content */
    font-family: SimHei, Arial, sans-serif;  /* Font supporting Chinese characters */
    background: white;
}
th, td {
    border: 1px solid black;  /* Add borders */
    padding: 8px;             /* Cell padding */
    text-align: center;       /* Center text */
}
th {
    background-color: #f2f2f2; /* Header background color */
}
"""

def parsing_image(image: Image, image_path: str) -> str:
    # Handle WMF images (PDFs)
    if image.ext == "wmf":
        image_path = image_path.replace(".wmf", ".jpg")
        if not pexists(image_path):
            wmf_to_images(image.blob, image_path)
    # Check for supported image types
    elif image.ext not in IMAGE_EXTENSIONS:
        raise ValueError(f"Unsupported image type {image.ext}")

    # Save image if it doesn't exist
    if not pexists(image_path):
        with open(image_path, "wb") as f:
            f.write(image.blob)
    return image_path


@tenacity_decorator
def wmf_to_images(blob: bytes, filepath: str):
    if not filepath.endswith(".jpg"):
        raise ValueError("filepath must end with .jpg")
    dirname = os.path.dirname(filepath)
    basename = os.path.basename(filepath).removesuffix(".jpg")
    with tempfile.TemporaryDirectory() as temp_dir:
        with open(pjoin(temp_dir, f"{basename}.wmf"), "wb") as f:
            f.write(blob)
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "jpg",
            pjoin(temp_dir, f"{basename}.wmf"),
            "--outdir",
            dirname,
        ]
        subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)

    assert pexists(filepath), f"File {filepath} does not exist"


def parse_groupshape(groupshape: GroupShape) -> list[dict[str, Length]]:
    """
    Parse a group shape to get the bounds of its child shapes.

    Args:
        groupshape (GroupShape): The group shape to parse.

    Returns:
        List[Dict[str, Length]]: The bounds of the child shapes.

    Raises:
        AssertionError: If the input is not a GroupShape.
    """
    assert isinstance(groupshape, GroupShape), "Input must be a GroupShape"

    # Get group bounds
    group_top_left_x = groupshape.left
    group_top_left_y = groupshape.top
    group_width = groupshape.width
    group_height = groupshape.height

    # Get shape bounds
    shape_top_left_x = min([sp.left for sp in groupshape.shapes])
    shape_top_left_y = min([sp.top for sp in groupshape.shapes])
    shape_width = (
        max([sp.left + sp.width for sp in groupshape.shapes]) - shape_top_left_x
    )
    shape_height = (
        max([sp.top + sp.height for sp in groupshape.shapes]) - shape_top_left_y
    )

    # Calculate bounds for each shape in the group
    group_shape_xy = []
    for sp in groupshape.shapes:
        group_shape_left = (
            sp.left - shape_top_left_x
        ) * group_width / shape_width + group_top_left_x
        group_shape_top = (
            sp.top - shape_top_left_y
        ) * group_height / shape_height + group_top_left_y
        group_shape_width = sp.width * group_width / shape_width
        group_shape_height = sp.height * group_height / shape_height

        group_shape_xy.append(
            {
                "left": Length(group_shape_left),
                "top": Length(group_shape_top),
                "width": Length(group_shape_width),
                "height": Length(group_shape_height),
            }
        )

    return group_shape_xy


def is_primitive(obj: Any) -> bool:
    """
    Check if an object is a primitive type or a collection of primitive types.

    Args:
        obj (Any): The object to check.

    Returns:
        bool: True if the object is a primitive type or a collection of primitive types, False otherwise.
    """
    if isinstance(obj, (list, tuple, set, frozenset)):
        return all(is_primitive(item) for item in obj)

    return isinstance(
        obj, (int, float, complex, bool, str, bytes, bytearray, type(None))
    )


DEFAULT_EXCLUDE: set[str] = {"element", "language_id", "ln", "placeholder_format"}


def dict_to_object(
    dict_obj: dict[str, Any], obj: Any, exclude: Optional[set[str]] = None
) -> None:
    """
    Apply dictionary values to an object.

    Args:
        dict_obj (Dict[str, Any]): The dictionary with values to apply.
        obj (Any): The object to apply values to.
        exclude (Optional[Set[str]]): The keys to exclude.
    """
    if exclude is None:
        exclude = set()

    for key, value in dict_obj.items():
        if key not in exclude and value is not None:
            setattr(obj, key, value)


def package_join(*paths: str) -> str:
    """
    Join paths with the appropriate separator for the platform.

    Args:
        *paths: The paths to join.

    Returns:
        str: The joined path.
    """
    _dir = pdirname(__file__)
    return pjoin(_dir, *paths)


class Config:
    """
    Configuration class for the application.
    """

    def __init__(
        self,
        rundir: Optional[str] = None,
        session_id: Optional[str] = None,
    ):
        """
        Initialize the configuration.

        Args:
            rundir (Optional[str]): The run directory.
            session_id (Optional[str]): The session ID.
            debug (bool): Whether to enable debug mode.
        """
        if rundir is not None:
            self.set_rundir(rundir)
        elif session_id is not None:
            self.set_session(session_id)
        else:
            raise ValueError("No session ID or run directory provided")

    def set_session(self, session_id: str) -> None:
        """
        Set the session ID and update the run directory.

        Args:
            session_id (str): The session ID.
        """
        self.session_id = session_id
        self.set_rundir(f"./runs/{session_id}")

    def set_rundir(self, rundir: str) -> None:
        """
        Set the run directory and create necessary subdirectories.

        Args:
            rundir (str): The run directory.
        """
        self.RUN_DIR = rundir
        self.IMAGE_DIR = pjoin(self.RUN_DIR, "images")

        for the_dir in [self.RUN_DIR, self.IMAGE_DIR]:
            os.makedirs(the_dir, exist_ok=True)

    def set_debug(self, debug: bool) -> None:
        """
        Set the debug mode.

        Args:
            debug (bool): Whether to enable debug mode.
        """
        self.DEBUG = debug

    def remove_rundir(self) -> None:
        """
        Remove the run directory and its subdirectories.
        """
        if pexists(self.RUN_DIR):
            shutil.rmtree(self.RUN_DIR)
        if pexists(self.IMAGE_DIR):
            shutil.rmtree(self.IMAGE_DIR)

    def __repr__(self) -> str:
        """
        Get a string representation of the configuration.

        Returns:
            str: A string representation of the configuration.
        """
        attrs = []
        for attr in dir(self):
            if not attr.startswith("_") and not callable(getattr(self, attr)):
                attrs.append(f"{attr}={getattr(self, attr)}")
        return f"Config({', '.join(attrs)})"


# Path utility functions
pjoin = os.path.join
pexists = os.path.exists
pbasename = os.path.basename
pdirname = os.path.dirname
