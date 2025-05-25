import tempfile
from copy import deepcopy
from dataclasses import asdict
from os.path import dirname, join

from dacite import from_dict
from pptx import Presentation as load_prs

from pptparse import Config, Presentation

TESTPATH = join(dirname(dirname(__file__)), "test.pptx")


def test_presentation():
    presentation = Presentation.from_file(TESTPATH, Config(tempfile.mkdtemp()))
    assert len(presentation.slides) == 1
    for sld in presentation.slides:
        sld.to_html(show_image=False)
    with tempfile.NamedTemporaryFile(suffix=".pptx") as tmp_file:
        attrs = asdict(deepcopy(presentation))
        from_dict(Presentation, attrs).save(tmp_file.name, layout_only=True)
        prs = load_prs(tmp_file.name)
        assert len(prs.slides) == 1
